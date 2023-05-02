# pip install python-pptx 

#ppt생성 
from pptx import Presentation # 라이브러리 
from pptx.util import Inches # 사진, 표등을 그리기 위해

prs = Presentation() # 파워포인트 객체 선언
prs.save(r'C:\Users\InSeong\Desktop\DT\test_ppt.pptx') #파워포인트 저장
#%%여러가지 슬라이드 추가 
prs = Presentation() # 파워포인트 객체 선언
for i in range(0, 11):
    title_slide_layout = prs.slide_layouts[i]  # 슬라이드 종류 선택
    prs.slides.add_slide(title_slide_layout) # 슬라이드 추가
prs.save(r'C:\Users\InSeong\Desktop\DT\test_ppt_slides11.pptx')


#%%
# placeholder: 슬라이드를 구성하는 요소 
prs = Presentation()
for i in range(0, 11):
    print("--------[%d] ------ "%(i))
    slide = prs.slides.add_slide(prs.slide_layouts[i])
    for shape in slide.placeholders:
        print(shape.name)
        print(shape.placeholder_format.idx)
        
#%%
# 본격적으로 꾸며보기 test.pptx
from pptx import Presentation # 라이브러리 

prs = Presentation() # 파워포인트 객체 선언
제목슬라이드= prs.slide_layouts[0] # 0 : 제목슬라이드에 해당
slide = prs.slides.add_slide(제목슬라이드) # 슬라이드 추가

#placeholder를 지정해줘야 편집하기 쉬움 
# 제목 - 제목에 값넣기
title = slide.placeholders[0] # 제목
title.text = "안녕하세요!" # 제목에 값 넣기

# 부제목
subtitle = slide.placeholders[1] # 제목상자는 placeholders[0], 부제목상자는 [1]
subtitle.text = "커리어하이 입니다 ! " 

# 저장
prs.save(r'C:\Users\InSeong\Desktop\DT\test_ppt.pptx')

#%% 단락추가하기 
#우선 두번째 슬라이드를 추가해줍니다. 
새로운슬라이드= prs.slide_layouts[1] # 새로운 슬라이드 레이아웃 지정
slide = prs.slides.add_slide(새로운슬라이드) # 슬라이드 추가

# 내용
body_shape = slide.placeholders[1]
tf = body_shape.text_frame #text frame을 쓰면 그안에 단락 편집 가능 
tf.text = '안녕하세요'

# 단락 추가
p = tf.add_paragraph()
p.text = '첫번째 주제가 들어가고 '
p.level = 1  # 1 : 들여쓰기 레벨

# 단락 추가 
p = tf.add_paragraph()
p.text = '그 안에 소주제는 여기들어가면 되겠네요 '
p.level = 2  # 2 : 들여쓰기 레벨

prs.save(r'C:\Users\InSeong\Desktop\DT\test_ppt.pptx')

#%% 그림정보넣기 

img_path = r'C:\Users\InSeong\Desktop\DT\test_img.png'

blank_slide_layout = prs.slide_layouts[6] # 6 : 제목/내용이 없는 '빈' 슬라이드
slide = prs.slides.add_slide(blank_slide_layout) #슬라이드 추가 

# 이미지 크기 조정 , width, hegith가 없을 경우 원본 사이즈로
left = top = Inches(1)
width = height = Inches(1)
# 이미지 삽입 with 이미지위치, 이미지사이즈 변수 
pic = slide.shapes.add_picture(img_path, left, top, width=width,height=height) 

left = Inches(3)
width = Inches(5)
height = Inches(4)
# 이미지 삽입 with 이미지위치,이미지사이즈 변수 
pic = slide.shapes.add_picture(img_path, left, top, width=width,height=height)

prs.save(r'C:\Users\InSeong\Desktop\DT\test_ppt.pptx')

#%% 표삽입 
from pptx.util import  Pt # 폰트조정
from pptx.dml.color import  RGBColor #컬러조정
title_only_slide_layout = prs.slide_layouts[5] 
slide = prs.slides.add_slide(title_only_slide_layout)

title_shape = slide.placeholders[0] 
title_shape.text = '커리어하이 주식운용부 주간손익 현황'
title_shape.text_frame.paragraphs[0].font.size = Pt(15)

rows = cols = 2
left = top = Inches(2)
width = Inches(6)
height = Inches(1)

#표 객체 삽입 with 표의 위치, 사이즈 , 표의 행과 열 개수정보
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# 표내부 열크기 조정
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)


# 표내부 글자 입력 및 글자폰트 수정 
table.cell(0, 0).text = '1팀 주간 손익'
table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(10)
table.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
table.cell(0, 0).text_frame.paragraphs[0].font.bold = True

table.cell(0, 1).text = '100,000,000'
table.cell(1, 0).text = '2팀 주간 손익'
table.cell(1, 1).text = '50,000,000'

prs.save(r'C:\Users\InSeong\Desktop\DT\test_ppt.pptx')



#%% 네이버 이메일 보내기 준비

import smtplib

smtp_info = dict({"smtp_server" : "smtp.naver.com", # SMTP 서버 주소
                  "smtp_user_id" : "pride0404@naver.com", #개인 이메일주소 xxx@naver.com
                  "smtp_user_pw" : "p4613404@" , #naver 비밀번호
                  "smtp_port" : 587}) # SMTP 서버 포트
#%% 네이버 이메일 보내기 
from email.mime.text import MIMEText 

메일내용 = "Python강의를 시작합니다."

# 메일 객체 생성 :메일 내용과 함께, 인코딩방식인 UTF-8을 명시해줍니다.
msg = MIMEText(_text = 메일내용, _charset = "utf-8") # 메일 내용
# MIMEText 객체는 아래와 같이 dictionary처럼 성분 추가가 가능합니다.
msg['Subject'] = "커리어하이에서 안내드립니다"    # 메일 제목
msg['From'] = "pride0404@naver.com"       # 송신자
msg['To'] = "pride0404@naver.com"     # 수신자

server = smtplib.SMTP(smtp_info["smtp_server"], smtp_info["smtp_port"])
# TLS 보안 연결
server.starttls()
# 로그인
server.login(smtp_info["smtp_user_id"], smtp_info["smtp_user_pw"])

# 서버에 이메일전송 , 메시지를 보낼때는 .as_string() 메소드를 사용해서 문자열로 바꿔줍니다
server.sendmail(msg['from'], msg['to'], msg.as_string())


#%% 네이버 이메일 보내기 (첨부파일 포함)
# 이메일 메시지에 다양한 형식을 누적해서 담기위함
from email.mime.multipart import MIMEMultipart
# 이메일 메시지 인코딩하기위함
from email import encoders
# 첨부파일을 담기위함 
from email.mime.base import MIMEBase
multi = MIMEMultipart(_subtype='mixed') #최종적으로 누적해서 담을 MIME객체생성
#파일(이미지)을 읽어서 MIMEbase에 담는다. 이후 multi 객체에 누적시킨다.
with open(r"C:\Users\InSeong\Desktop\DT\test_img.png", 'rb') as fp:
    msg = MIMEBase('application',  _subtype='octect-stream')
    msg.set_payload(fp.read())
    encoders.encode_base64(msg)
#보낼 파일명에 확장자를 붙혀서 수신자의 조회가 용이하게 합니다.
# 'Content-Disposition', 'attachment' : 첨부파일을 첨부하게 하는 명령어 
msg.add_header('Content-Disposition', 'attachment', filename='커리어하이이미지.png')
multi.attach(msg)
#두번째 파일(ppt)을 읽어서 MIMEbase에 담는다. 이후 multi 객체에 누적시킨다.
with open(r"C:\Users\InSeong\Desktop\DT\test_ppt.pptx", 'rb') as fp:
    msg = MIMEBase('application',  _subtype='octect-stream')
    msg.set_payload(fp.read())
    encoders.encode_base64(msg)
msg.add_header('Content-Disposition', 'attachment', filename='pptx로만든ppt.pptx')
multi.attach(msg)
#이메일 내용 텍스트를 담은 MIMEText도 multi 객체에 누적시킨다.
메일내용 = "이미지파일 보내드립니다."
multi.attach(MIMEText(_text = 메일내용, _charset = "utf-8") )
#이하는 위와 같이 MIME객체를 가지고 송수신자를 결정하여 메일을 보낸다.
multi['subject'] = "커리어하이에서 보내드립니다." 
multi['from'] ="pride0404@naver.com"
multi['to'] = "pride0404@naver.com" 
server = smtplib.SMTP(smtp_info["smtp_server"], smtp_info["smtp_port"])
server.starttls()
server.login(smtp_info["smtp_user_id"], smtp_info["smtp_user_pw"])
server.sendmail(multi['from'], multi['to'], multi.as_string())
